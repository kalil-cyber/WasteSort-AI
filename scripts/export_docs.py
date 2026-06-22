"""
export_docs.py
Génère les livrables PPTX et PDF pour WasteSort AI Suite.
"""

from pathlib import Path
import re
import textwrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    Image as PdfImage,
    PageBreak,
    Paragraph,
    Preformatted,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parent.parent
EXPORTS = ROOT / "exports"
LOGO = ROOT / "assets" / "logo.svg"
EMSI_LOGO = ROOT / "assets" / "emsi-logo-official.png"
REPORT = ROOT / "report.md"
GENERATED_ASSETS = ROOT / "assets"
PROJECT_INFO = {
    "Nom de l'étudiant": "......",
    "Nom de la prof": "......",
    "Classe / formation": "......",
    "Établissement": "EMSI / ......",
    "Année scolaire": "......",
}
VISUALS = {
    0: GENERATED_ASSETS / "wastesort-slide-hero.png",
    4: GENERATED_ASSETS / "wastesort-slide-categories.png",
    6: GENERATED_ASSETS / "wastesort-slide-ai-model.png",
    7: GENERATED_ASSETS / "wastesort-slide-ai-model.png",
    8: GENERATED_ASSETS / "wastesort-slide-hero.png",
}


SLIDES = [
    (
        "WasteSort AI Suite",
        [
            "Plateforme IA multi-modèles pour le tri et l'écologie",
            "MLP - CNN/MobileNetV2 - LSTM - Seq2Seq",
            "Application web Streamlit",
            "Étudiant : ......",
            "Prof : ......",
        ],
    ),
    (
        "Pourquoi ce projet ?",
        [
            "Montrer plusieurs familles de modèles IA dans un seul projet.",
            "Aider un campus à gérer le tri, l'énergie et les questions utilisateurs.",
            "Rendre l'IA compréhensible avec une interface simple.",
        ],
    ),
    (
        "Problématique",
        [
            "Comment regrouper plusieurs modèles IA dans une seule plateforme ?",
            "Comment traiter images, tableaux, séries temporelles et texte ?",
        ],
    ),
    (
        "Objectifs",
        [
            "Reconnaître les déchets par image.",
            "Calculer un score écologique.",
            "Prévoir une consommation énergétique.",
            "Répondre aux questions sur le tri.",
            "Comparer MLP, CNN, LSTM et Seq2Seq.",
        ],
    ),
    (
        "Modules De La Plateforme",
        [
            "Déchets : CNN / MobileNetV2",
            "Score écologique : MLP",
            "Prévision énergie : RNN / LSTM",
            "Assistant tri : Seq2Seq / FAQ",
            "Dashboard : Streamlit",
        ],
    ),
    (
        "Technologies Utilisées",
        [
            "Python",
            "TensorFlow / Keras",
            "MobileNetV2 en Transfer Learning",
            "MLP, RNN/LSTM et Seq2Seq",
            "Streamlit",
            "Scikit-learn, Matplotlib et Pillow",
            "Dataset Kaggle Garbage Classification",
        ],
    ),
    (
        "Méthodologie",
        [
            "Module image : MobileNetV2 pour les déchets.",
            "Module MLP : score écologique sur données tabulaires.",
            "Module LSTM : prévision sur séries temporelles.",
            "Module Seq2Seq : assistant de tri.",
            "Intégration dans une interface Streamlit.",
        ],
    ),
    (
        "Module Image",
        [
            "Reconnaît les déchets à partir d'une photo.",
            "Utilise MobileNetV2 avec Transfer Learning.",
            "Affiche confiance, probabilités et statut de fiabilité.",
            "Accuracy validation : environ 81,5 %.",
        ],
    ),
    (
        "Comparaison Des Modèles",
        [
            "MLP : score écologique à partir de données tabulaires.",
            "CNN / MobileNetV2 : classification d'images.",
            "RNN / LSTM : prévision d'une série temporelle.",
            "Seq2Seq : assistant de questions/réponses.",
            "Chaque modèle correspond à un type de donnée.",
        ],
    ),
    (
        "Interface Utilisateur",
        [
            "Navigation entre les modules.",
            "Test image, score, prévision et chatbot.",
            "Résultats affichés simplement.",
            "Interface adaptée à une démonstration publique.",
        ],
    ),
    (
        "Résultats",
        [
            "MobileNetV2 : environ 81,5 % en validation.",
            "MLP : score écologique explicable.",
            "LSTM : prévision basée sur l'historique.",
            "Seq2Seq/FAQ : réponses simples sur le tri.",
        ],
    ),
    (
        "Limites",
        [
            "Les modules MLP, LSTM et Seq2Seq sont des prototypes.",
            "Le modèle image dépend de la qualité de la photo.",
            "Certaines classes de déchets se ressemblent.",
            "Des datasets réels amélioreraient la plateforme.",
        ],
    ),
    (
        "Améliorations Possibles",
        [
            "Entraîner un vrai MLP sur des données campus.",
            "Entraîner un vrai LSTM sur consommation réelle.",
            "Créer un vrai chatbot Seq2Seq ou Transformer.",
            "Ajouter une application mobile.",
            "Connecter la plateforme à une base de données.",
        ],
    ),
    (
        "Conclusion",
        [
            "WasteSort AI Suite regroupe plusieurs familles de modèles.",
            "Chaque module illustre un type de donnée différent.",
            "La plateforme est claire, visuelle et pédagogique.",
            "Le projet est adapté à une présentation académique.",
        ],
    ),
]


def clean_markdown(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = text.replace("`", "")
    return text


def add_title(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(15, 81, 50)


def add_footer(slide):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(11.8), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = "WasteSort AI Suite - MLP - CNN/MobileNetV2 - LSTM - Seq2Seq"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.RIGHT


def add_ppt_logo(slide, left, top, size):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        size,
        size,
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(232, 245, 233)
    box.line.width = Pt(2)

    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left + size * 0.18,
        top + size * 0.16,
        size * 0.64,
        size * 0.64,
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(25, 135, 84)
    circle.line.color.rgb = RGBColor(25, 135, 84)

    leaf = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left + size * 0.28,
        top + size * 0.24,
        size * 0.44,
        size * 0.44,
    )
    leaf.fill.solid()
    leaf.fill.fore_color.rgb = RGBColor(116, 184, 74)
    leaf.line.color.rgb = RGBColor(116, 184, 74)

    text = slide.shapes.add_textbox(left + size * 0.20, top + size * 0.70, size * 0.60, size * 0.22)
    p = text.text_frame.paragraphs[0]
    p.text = "WS"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(15, 81, 50)


def add_emsi_badge(slide, left, top, width, height):
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
    badge.line.color.rgb = RGBColor(15, 81, 50)
    badge.line.width = Pt(2)

    icon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left + Inches(0.18),
        top + Inches(0.14),
        Inches(0.52),
        Inches(0.52),
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = RGBColor(15, 81, 50)
    icon.line.color.rgb = RGBColor(15, 81, 50)

    icon_text = slide.shapes.add_textbox(left + Inches(0.26), top + Inches(0.25), Inches(0.36), Inches(0.2))
    p = icon_text.text_frame.paragraphs[0]
    p.text = "E"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)

    text = slide.shapes.add_textbox(left + Inches(0.82), top + Inches(0.15), width - Inches(1.0), height - Inches(0.2))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = "EMSI"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(15, 81, 50)
    p2 = tf.add_paragraph()
    p2.text = "École Marocaine des Sciences de l'Ingénieur"
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(71, 85, 105)


def add_visual(slide, visual_path: Path, left, top, width):
    if visual_path.exists():
        picture = slide.shapes.add_picture(str(visual_path), left, top, width=width)
        picture.line.color.rgb = RGBColor(226, 232, 240)
        return picture
    return None


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, (title, bullets) in enumerate(SLIDES):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)

        if index == 0:
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.85), Inches(5.9), Inches(1.0))
            p = title_box.text_frame.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(46)
            p.font.color.rgb = RGBColor(15, 81, 50)

            subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(5.7), Inches(2.2))
            tf = subtitle_box.text_frame
            for i, item in enumerate(bullets):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = item
                para.font.size = Pt(21 if i == 0 else 16)
                para.font.color.rgb = RGBColor(51, 65, 85)
                para.space_after = Pt(7)

            info_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.35), Inches(5.7), Inches(1.45))
            info_tf = info_box.text_frame
            for i, (label, value) in enumerate(PROJECT_INFO.items()):
                para = info_tf.paragraphs[0] if i == 0 else info_tf.add_paragraph()
                para.text = f"{label} : {value}"
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(15, 81, 50)
                para.space_after = Pt(2)

            if EMSI_LOGO.exists():
                slide.shapes.add_picture(
                    str(EMSI_LOGO),
                    Inches(0.8),
                    Inches(5.78),
                    Inches(4.2),
                    Inches(1.15),
                )
            else:
                add_emsi_badge(slide, Inches(0.8), Inches(6.02), Inches(3.9), Inches(0.78))

            visual = VISUALS.get(index)
            if visual and visual.exists():
                add_visual(slide, visual, Inches(6.55), Inches(0.65), Inches(6.0))
            else:
                add_ppt_logo(slide, Inches(9.0), Inches(1.25), Inches(2.45))
        else:
            add_title(slide, title)
            visual = VISUALS.get(index)
            has_visual = visual is not None and visual.exists()
            content_width = Inches(5.45) if has_visual else Inches(11.0)
            content = slide.shapes.add_textbox(Inches(1.0), Inches(1.45), content_width, Inches(4.9))
            tf = content.text_frame
            tf.word_wrap = True
            for i, item in enumerate(bullets):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = item
                para.level = 0
                para.font.size = Pt(22)
                para.font.color.rgb = RGBColor(30, 41, 59)
                para.space_after = Pt(12)
                para.text = f"• {item}"

            if has_visual:
                add_visual(slide, visual, Inches(6.85), Inches(1.45), Inches(5.45))

        add_footer(slide)

    output = EXPORTS / "WasteSort_AI_Suite_Presentation.pptx"
    prs.save(output)
    return output


def markdown_blocks(path: Path):
    lines = path.read_text(encoding="utf-8").splitlines()
    blocks = []
    buffer = []
    in_code = False
    for line in lines:
        if line.startswith("```"):
            in_code = not in_code
            if in_code and buffer:
                blocks.append("\n".join(buffer))
                buffer = []
            continue
        if in_code:
            buffer.append(line)
            continue
        if not line.strip():
            if buffer:
                blocks.append("\n".join(buffer))
                buffer = []
            continue
        buffer.append(line)
    if buffer:
        blocks.append("\n".join(buffer))
    return blocks


def build_pdf():
    output = EXPORTS / "WasteSort_AI_Suite_Rapport.pdf"
    doc = SimpleDocTemplate(
        str(output),
        pagesize=A4,
        rightMargin=2 * cm,
        leftMargin=2 * cm,
        topMargin=1.7 * cm,
        bottomMargin=1.7 * cm,
    )
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="DocTitle",
            parent=styles["Title"],
            alignment=TA_CENTER,
            textColor=colors.HexColor("#0F5132"),
            fontSize=22,
            leading=28,
            spaceAfter=18,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SectionTitle",
            parent=styles["Heading2"],
            textColor=colors.HexColor("#198754"),
            fontSize=15,
            leading=20,
            spaceBefore=12,
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="BodyTextClean",
            parent=styles["BodyText"],
            alignment=TA_LEFT,
            fontSize=10.5,
            leading=15,
            spaceAfter=7,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Caption",
            parent=styles["BodyText"],
            fontSize=9,
            leading=12,
            textColor=colors.HexColor("#475569"),
            alignment=TA_CENTER,
            spaceAfter=10,
        )
    )
    code_style = ParagraphStyle(
        name="CodeBlock",
        fontName="Courier",
        fontSize=9,
        leading=12,
        leftIndent=12,
        rightIndent=12,
        textColor=colors.HexColor("#0F172A"),
        backColor=colors.HexColor("#F1F5F9"),
        borderPadding=8,
        spaceBefore=6,
        spaceAfter=10,
    )

    story = []
    story.append(Paragraph("Rapport Professionnel - WasteSort AI Suite", styles["DocTitle"]))
    story.append(Paragraph("Plateforme IA multi-modèles pour le tri et l'écologie", styles["BodyTextClean"]))
    story.append(Spacer(1, 0.25 * cm))

    info_rows = [["Information", "Détail"]]
    info_rows.extend([[label, value] for label, value in PROJECT_INFO.items()])
    info_rows.append(["Technologies principales", "Python, TensorFlow/Keras, MobileNetV2, Streamlit"])
    info_table = Table(info_rows, colWidths=[6 * cm, 9 * cm], hAlign="LEFT")
    info_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8F5E9")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0F5132")),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#CBD5E1")),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    story.append(info_table)
    story.append(Spacer(1, 0.3 * cm))

    for block in markdown_blocks(REPORT):
        block = clean_markdown(block)
        if block.startswith("# "):
            continue
        if block.startswith("## "):
            story.append(Paragraph(block[3:], styles["SectionTitle"]))
            continue
        if block.startswith("### "):
            story.append(Paragraph(block[4:], styles["Heading3"]))
            continue
        if block.startswith("!["):
            match = re.match(r"!\[(.*?)\]\((.*?)\)", block)
            if match:
                alt_text, image_rel_path = match.groups()
                image_path = ROOT / image_rel_path
                if image_path.exists():
                    if image_path.suffix.lower() == ".svg":
                        badge = Table(
                            [["EMSI", "École Marocaine des Sciences de l'Ingénieur"]],
                            colWidths=[3 * cm, 12 * cm],
                            hAlign="CENTER",
                        )
                        badge.setStyle(
                            TableStyle(
                                [
                                    ("BACKGROUND", (0, 0), (0, 0), colors.HexColor("#0F5132")),
                                    ("TEXTCOLOR", (0, 0), (0, 0), colors.white),
                                    ("TEXTCOLOR", (1, 0), (1, 0), colors.HexColor("#0F5132")),
                                    ("FONTNAME", (0, 0), (-1, -1), "Helvetica-Bold"),
                                    ("FONTSIZE", (0, 0), (0, 0), 16),
                                    ("FONTSIZE", (1, 0), (1, 0), 11),
                                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#0F5132")),
                                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                                    ("TOPPADDING", (0, 0), (-1, -1), 10),
                                    ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
                                ]
                            )
                        )
                        story.append(Spacer(1, 0.15 * cm))
                        story.append(badge)
                        story.append(Paragraph(alt_text, styles["Caption"]))
                        continue
                    story.append(Spacer(1, 0.15 * cm))
                    story.append(PdfImage(str(image_path), width=15.5 * cm, height=8.7 * cm))
                    story.append(Paragraph(alt_text, styles["Caption"]))
            continue
        if "↓" in block or block.startswith("Photo utilisateur") or block.startswith("Dataset Kaggle"):
            story.append(Preformatted(block, code_style))
            continue
        if block.startswith("|"):
            rows = [row.strip("|").split("|") for row in block.splitlines() if row.startswith("|")]
            rows = [[cell.strip() for cell in row] for row in rows if "---" not in "".join(row)]
            if rows:
                table = Table(rows, hAlign="LEFT")
                table.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8F5E9")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0F5132")),
                            ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#CBD5E1")),
                            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                            ("VALIGN", (0, 0), (-1, -1), "TOP"),
                            ("LEFTPADDING", (0, 0), (-1, -1), 6),
                            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                        ]
                    )
                )
                story.append(table)
                story.append(Spacer(1, 0.25 * cm))
            continue

        for paragraph in block.splitlines():
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            if paragraph.startswith("- "):
                paragraph = "• " + paragraph[2:]
            if re.match(r"^\d+\. ", paragraph):
                paragraph = paragraph
            wrapped = "<br/>".join(textwrap.wrap(paragraph, width=105))
            story.append(Paragraph(wrapped, styles["BodyTextClean"]))

    story.append(PageBreak())
    story.append(Paragraph("Annexe - Résumé Pour Présentation", styles["DocTitle"]))
    story.append(
        Paragraph(
            "WasteSort AI Suite regroupe plusieurs familles de modèles : MLP pour les données tabulaires, "
            "CNN/MobileNetV2 pour les images, LSTM pour les séries temporelles et Seq2Seq pour l'assistant de tri.",
            styles["BodyTextClean"],
        )
    )
    doc.build(story)
    return output


def main():
    EXPORTS.mkdir(exist_ok=True)
    pptx_path = build_pptx()
    pdf_path = build_pdf()
    print(f"PPTX généré : {pptx_path}")
    print(f"PDF généré : {pdf_path}")


if __name__ == "__main__":
    main()
